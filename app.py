from io import BytesIO
from docx import Document
from flask import Flask, render_template, request, send_file
from pptx import Presentation

app = Flask(__name__)

def save_text_to_docx(text, docx_filename):
    doc = Document()
    doc.add_paragraph(text)
    doc.save(docx_filename)
    
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    ppt_file = request.files['pptFile']
    if not ppt_file:
        return 'No file uploaded', 400

    try:
        def extract_text_from_ppt(ppt_file):
            ppt = Presentation(ppt_file)
            extracted_text = ""

            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                extracted_text += run.text + " "
            
            return extracted_text

        extracted_text = extract_text_from_ppt(ppt_file)
        output = save_text_to_docx(extracted_text, 'converted.docx')
    except Exception as e:
        return f'Error converting file: {str(e)}', 500

    return send_file(output, mimetype='application/msword', as_attachment=True, download_name='converted.docx')

@app.route('/download')
def download():
    return send_file('sample.docx', as_attachment=True, download_name='sample.docx')

if __name__ == '__main__':
    app.run()
